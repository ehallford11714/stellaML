"""Unstructured data extraction and NLP utilities (HTML/XML/PPTX/URLs)."""

from __future__ import annotations

import re
import xml.etree.ElementTree as ET
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any


@dataclass(slots=True)
class UnstructuredDocument:
    source: str
    text: str
    metadata: dict[str, Any]


@dataclass(slots=True)
class TextEDAReport:
    num_chars: int
    num_tokens: int
    num_sentences: int
    top_terms: list[tuple[str, int]]
    vocabulary_size: int


def fetch_url_with_requests(url: str, timeout: int = 30) -> str:
    import requests

    response = requests.get(url, timeout=timeout)
    response.raise_for_status()
    return response.text


def fetch_url_with_selenium(url: str, wait_seconds: int = 3) -> str:
    """Fetch rendered HTML using Selenium (headless Chrome expected)."""
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options

    options = Options()
    options.add_argument("--headless=new")
    options.add_argument("--no-sandbox")
    options.add_argument("--disable-dev-shm-usage")

    driver = webdriver.Chrome(options=options)
    try:
        driver.get(url)
        driver.implicitly_wait(wait_seconds)
        return driver.page_source
    finally:
        driver.quit()


def parse_html_text(html: str) -> str:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return " ".join(soup.stripped_strings)
    except Exception:
        no_tags = re.sub(r"<[^>]+>", " ", html)
        return " ".join(no_tags.split())


def parse_xml_text(xml_content: str) -> str:
    root = ET.fromstring(xml_content)
    return " ".join(text.strip() for text in root.itertext() if text and text.strip())


def parse_pptx_text(pptx_path: str | Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def extract_unstructured(
    source: str,
    source_type: str = "url",
    fetch_mode: str = "requests",
) -> UnstructuredDocument:
    """Extract text from URL/HTML/XML/PPTX into normalized document object."""
    source_type = source_type.lower()

    if source_type == "url":
        html = fetch_url_with_selenium(source) if fetch_mode == "selenium" else fetch_url_with_requests(source)
        text = parse_html_text(html)
        return UnstructuredDocument(source=source, text=text, metadata={"source_type": "url", "fetch_mode": fetch_mode})

    if source_type == "html":
        return UnstructuredDocument(source="inline_html", text=parse_html_text(source), metadata={"source_type": "html"})

    if source_type == "xml":
        return UnstructuredDocument(source="inline_xml", text=parse_xml_text(source), metadata={"source_type": "xml"})

    if source_type == "pptx":
        path = Path(source)
        return UnstructuredDocument(source=str(path), text=parse_pptx_text(path), metadata={"source_type": "pptx"})

    raise ValueError(f"Unsupported source_type: {source_type}")


def autoeda_text(text: str, top_k: int = 20) -> TextEDAReport:
    tokens = tokenize_text(text)
    sentences = [s for s in re.split(r"[.!?]+", text) if s.strip()]
    counts = Counter(tokens)
    return TextEDAReport(
        num_chars=len(text),
        num_tokens=len(tokens),
        num_sentences=len(sentences),
        top_terms=counts.most_common(top_k),
        vocabulary_size=len(counts),
    )


def tokenize_text(text: str) -> list[str]:
    return re.findall(r"[A-Za-z0-9']+", text.lower())


def ngrams(text: str, n: int = 2, top_k: int = 20) -> list[tuple[str, int]]:
    toks = tokenize_text(text)
    grams = [" ".join(toks[i : i + n]) for i in range(max(0, len(toks) - n + 1))]
    return Counter(grams).most_common(top_k)


def nltk_ngrams(text: str, n: int = 2, top_k: int = 20) -> list[tuple[str, int]]:
    import nltk

    toks = tokenize_text(text)
    grams = [" ".join(g) for g in nltk.ngrams(toks, n)]
    return Counter(grams).most_common(top_k)


def spacy_entities(text: str, model: str = "en_core_web_sm") -> list[tuple[str, str]]:
    try:
        import spacy
        nlp = spacy.load(model)
    except Exception:
        return []
    doc = nlp(text)
    return [(ent.text, ent.label_) for ent in doc.ents]


def run_unstructured_nlp_pipeline(
    source: str,
    source_type: str = "url",
    fetch_mode: str = "requests",
    ngram_n: int = 2,
) -> dict[str, Any]:
    document = extract_unstructured(source, source_type=source_type, fetch_mode=fetch_mode)
    report = autoeda_text(document.text)
    return {
        "document": document,
        "text_eda": report,
        "ngrams": ngrams(document.text, n=ngram_n),
        "entities": spacy_entities(document.text),
    }
